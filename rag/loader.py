import os
import logging
from typing import List, Dict, Any, Generator
import pdfplumber
from pptx import Presentation

logger = logging.getLogger(__name__)

class DocumentLoader:
    def scan_course_directory(self, base_path: str) -> Generator[Dict[str, Any], None, None]:
        """
        Recursively scans the base directory for course materials.
        Yields a dictionary containing file metadata and extracted text.
        Structure:
        {
            "course_name": str,
            "file_name": str,
            "file_path": str,
            "file_type": str,
            "content": List[Dict] # [{"text": "...", "page_num": 1}, ...]
        }
        """
        if not os.path.exists(base_path):
            logger.error(f"Directory not found: {base_path}")
            return

        # Iterate through immediate subdirectories (Courses)
        for entry in os.scandir(base_path):
            if entry.is_dir():
                course_name = entry.name
                # Walk through the course directory
                for root, _, files in os.walk(entry.path):
                    for file in files:
                        file_path = os.path.join(root, file)
                        file_ext = os.path.splitext(file)[1].lower()
                        
                        content = []
                        try:
                            if file_ext == '.pdf':
                                content = self._extract_pdf(file_path)
                            elif file_ext == '.pptx':
                                content = self._extract_pptx(file_path)
                            else:
                                continue # Skip unsupported files

                            if content:
                                yield {
                                    "course_name": course_name,
                                    "file_name": file,
                                    "file_path": file_path,
                                    "file_type": file_ext,
                                    "content": content
                                }
                        except Exception as e:
                            logger.exception("Error extracting %s", file_path)

    def _extract_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        extracted_pages = []
        try:
            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages)
                logger.info("PDF extract start: path=%s pages=%s", file_path, total_pages)
                for i, page in enumerate(pdf.pages):
                    text = ""
                    try:
                        text = page.extract_text() or ""
                    except Exception:
                        logger.exception("PDF extract_text failed: %s page=%s", file_path, i + 1)
                        text = ""

                    text = text.strip()
                    logger.info("PDF page=%s extract_text_len=%s", i + 1, len(text))

                    if not text:
                        try:
                            ocr_text = self._ocr_pdf_page(file_path=file_path, page_index=i) or ""
                            ocr_text = ocr_text.strip()
                            logger.info("PDF page=%s ocr_len=%s", i + 1, len(ocr_text))
                            text = ocr_text
                        except Exception:
                            logger.exception("PDF OCR failed: %s page=%s", file_path, i + 1)
                            text = ""

                    if text:
                        extracted_pages.append({"text": text, "page_num": i + 1})
        except Exception as e:
            logger.exception("PDF extraction error for %s", file_path)
            raise e

        if extracted_pages:
            logger.info("PDF extract done: path=%s pages_with_text=%s", file_path, len(extracted_pages))
        else:
            logger.info("PDF extract done: path=%s pages_with_text=0 (likely scanned or empty)", file_path)

        return extracted_pages

    def _ocr_pdf_page(self, file_path: str, page_index: int) -> str:
        try:
            import numpy as np
            import pypdfium2 as pdfium
            from paddleocr import PaddleOCR
        except Exception as e:
            raise RuntimeError(f"OCR dependencies not available: {e}")

        ocr = getattr(self, "_paddle_ocr", None)
        if ocr is None:
            ocr = PaddleOCR(use_angle_cls=True, lang="ch")
            setattr(self, "_paddle_ocr", ocr)

        pdf = pdfium.PdfDocument(file_path)
        try:
            if page_index < 0 or page_index >= len(pdf):
                return ""
            page = pdf[page_index]
            try:
                bitmap = page.render(scale=2.0)
                img = bitmap.to_pil()
            finally:
                try:
                    page.close()
                except Exception:
                    pass
        finally:
            try:
                pdf.close()
            except Exception:
                pass

        if img.mode != "RGB":
            img = img.convert("RGB")

        arr = np.array(img)
        result = ocr.ocr(arr, cls=True)
        texts: List[str] = []
        for item in result or []:
            if not item or len(item) < 2:
                continue
            txt = item[1][0] if item[1] else ""
            if txt:
                texts.append(str(txt))
        return "\n".join(texts).strip()

    def _extract_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        extracted_slides = []
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text.strip())
                
                full_text = "\n".join(text_parts).strip()
                if full_text:
                    extracted_slides.append({
                        "text": full_text,
                        "page_num": i + 1
                    })
        except Exception as e:
            logger.exception("PPTX extraction error for %s", file_path)
        return extracted_slides

loader = DocumentLoader()
